from pptx import Presentation
a = ['Module random.', 'Random variable generators.']
b = ['integers:', 'uniform within range  ']  
c = ['sequences: ', 'pick random element', 'generate random permutation ']
d = ['distributions on the real line:', 'uniform', 'triangular', 'lognormal']
e = ['General notes:', 'The period is 2**19937-1.', 'It is the most extensively tested generators']

prs = Presentation()
title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)

title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = a[0]
subtitle.text = a[1]

slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = b[0]
subtitle.text = b[1]

slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = c[0]
subtitle.text = c[1] + '\n' + c[2]


slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = d[0]
subtitle.text = d[1] + '\n' + d[2] + '\n' + d[3]

slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = e[0]
subtitle.text = e[1] + '\n' + e[2]

prs.save('prez.pptx')
